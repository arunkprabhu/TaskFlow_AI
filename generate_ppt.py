"""Generate an attractive hackathon PPT for MeetingToTask app."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import io
from PIL import Image, ImageDraw, ImageFont
import math

# ── Design tokens ──────────────────────────────────────────────────────────────
DARK_BG      = RGBColor(0x0D, 0x1B, 0x2A)   # deep navy
ACCENT       = RGBColor(0x00, 0xD4, 0xFF)   # electric cyan
ACCENT2      = RGBColor(0x7B, 0x2F, 0xFF)   # vivid purple
ACCENT3      = RGBColor(0x00, 0xFF, 0xC8)   # mint green
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY   = RGBColor(0xCC, 0xD6, 0xE0)
CARD_BG      = RGBColor(0x16, 0x2A, 0x3E)
YELLOW       = RGBColor(0xFF, 0xD6, 0x00)

W, H = Inches(13.33), Inches(7.5)   # 16:9 widescreen

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]  # completely blank

# ── Helpers ────────────────────────────────────────────────────────────────────

def bg(slide, color=DARK_BG):
    """Fill slide background."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def box(slide, l, t, w, h, color, alpha=None, radius=None):
    """Add a filled rectangle."""
    shape = slide.shapes.add_shape(1, l, t, w, h)   # MSO_SHAPE_TYPE.RECTANGLE = 1
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def txt(slide, text, l, t, w, h, size=24, bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, italic=False, wrap=True):
    tf = slide.shapes.add_textbox(l, t, w, h)
    tf.word_wrap = wrap
    p  = tf.text_frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size   = Pt(size)
    run.font.bold   = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tf

def gradient_image(w_px, h_px, c1, c2, vertical=True):
    """Return PIL image with a linear gradient."""
    img = Image.new("RGB", (w_px, h_px))
    draw = ImageDraw.Draw(img)
    for i in range(h_px if vertical else w_px):
        t = i / (h_px if vertical else w_px)
        r = int(c1[0] + (c2[0]-c1[0])*t)
        g = int(c1[1] + (c2[1]-c1[1])*t)
        b = int(c1[2] + (c2[2]-c1[2])*t)
        if vertical:
            draw.line([(0,i),(w_px,i)], fill=(r,g,b))
        else:
            draw.line([(i,0),(i,h_px)], fill=(r,g,b))
    return img

def add_pil_image(slide, pil_img, l, t, w, h):
    buf = io.BytesIO()
    pil_img.save(buf, format="PNG")
    buf.seek(0)
    slide.shapes.add_picture(buf, l, t, w, h)

def circuit_image(w_px, h_px):
    """Draw a techy circuit-board pattern."""
    img = Image.new("RGB", (w_px, h_px), (13, 27, 42))
    draw = ImageDraw.Draw(img)
    color = (0, 212, 255, 80)
    # horizontal & vertical lines
    import random; random.seed(42)
    for _ in range(30):
        x = random.randint(0, w_px)
        y = random.randint(0, h_px)
        l2 = random.randint(40, 200)
        if random.random() > 0.5:
            draw.line([(x,y),(x+l2,y)], fill=(0,80,120), width=1)
            draw.ellipse([x+l2-3, y-3, x+l2+3, y+3], fill=(0,212,255))
        else:
            draw.line([(x,y),(x,y+l2)], fill=(0,80,120), width=1)
            draw.ellipse([x-3, y+l2-3, x+3, y+l2+3], fill=(0,212,255))
    # dots grid
    for xi in range(0, w_px, 60):
        for yi in range(0, h_px, 60):
            draw.ellipse([xi-2,yi-2,xi+2,yi+2], fill=(0,60,90))
    return img

def wave_image(w_px, h_px):
    """Sine-wave decoration."""
    img = Image.new("RGBA", (w_px, h_px), (0,0,0,0))
    draw = ImageDraw.Draw(img)
    for wave in range(3):
        amp   = 20 - wave*5
        freq  = 0.015 + wave*0.005
        phase = wave * 1.2
        col   = [(0,212,255,120),(123,47,255,100),(0,255,200,80)][wave]
        pts   = []
        for x in range(w_px+1):
            y = int(h_px/2 + amp*math.sin(freq*x + phase))
            pts.append((x,y))
        draw.line(pts, fill=col, width=2)
    return img

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Hero / Title
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
bg(slide)

# gradient left panel
grad = gradient_image(700, 540, (0,60,120), (13,27,42), vertical=False)
add_pil_image(slide, grad, Inches(0), Inches(0), Inches(6.5), Inches(7.5))

# circuit art on right
circ = circuit_image(840, 540)
add_pil_image(slide, circ, Inches(6.3), Inches(0), Inches(7), Inches(7.5))

# accent top bar
box(slide, Inches(0), Inches(0), W, Inches(0.08), ACCENT)

# glowing accent side strip
box(slide, Inches(6.3), Inches(0), Inches(0.06), H, ACCENT)

# hackathon badge
badge = box(slide, Inches(0.4), Inches(0.3), Inches(2.4), Inches(0.5), ACCENT2)
txt(slide, "🏆  HACKATHON 2026", Inches(0.42), Inches(0.29), Inches(2.4), Inches(0.52),
    size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# main title
txt(slide, "MeetingToTask", Inches(0.4), Inches(1.1), Inches(5.8), Inches(1.3),
    size=58, bold=True, color=ACCENT, align=PP_ALIGN.LEFT)

# subtitle
txt(slide, "Turn Any Meeting into Actionable Tasks — Instantly.",
    Inches(0.4), Inches(2.3), Inches(5.7), Inches(0.9),
    size=22, bold=False, color=WHITE, align=PP_ALIGN.LEFT)

# tagline with accent
txt(slide, "AI-Powered  ·  Private  ·  Zero Manual Work",
    Inches(0.4), Inches(3.15), Inches(5.7), Inches(0.6),
    size=16, color=ACCENT3, align=PP_ALIGN.LEFT)

# divider line
box(slide, Inches(0.4), Inches(3.75), Inches(3.5), Inches(0.04), ACCENT)

# tech pills
for i, (label, color) in enumerate([
        ("React + TypeScript", ACCENT),
        ("FastAPI Python", ACCENT3),
        ("Ollama Local AI", ACCENT2),
        ("Monday.com", YELLOW)]):
    pill = box(slide, Inches(0.4 + i*1.38), Inches(4.0), Inches(1.28), Inches(0.38), color)
    txt(slide, label, Inches(0.4 + i*1.38), Inches(3.99), Inches(1.28), Inches(0.4),
        size=10, bold=True, color=DARK_BG, align=PP_ALIGN.CENTER)

# bottom wave
wave = wave_image(960, 80)
add_pil_image(slide, wave, Inches(0), Inches(6.9), W, Inches(0.6))

txt(slide, "Paste notes ➜ AI extracts tasks ➜ Push to Monday.com",
    Inches(0.4), Inches(5.1), Inches(5.8), Inches(0.6),
    size=15, italic=True, color=LIGHT_GRAY, align=PP_ALIGN.LEFT)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — The Problem
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
bg(slide)
box(slide, Inches(0), Inches(0), W, Inches(0.07), ACCENT2)

txt(slide, "😫  The Problem", Inches(0.5), Inches(0.2), Inches(8), Inches(0.8),
    size=38, bold=True, color=ACCENT2)

problems = [
    ("⏰", "Meetings end — action items get forgotten"),
    ("📝", "Someone has to manually rewrite notes into tasks"),
    ("🔁", "Copy-paste from notes → project board wastes hours"),
    ("❓", "No clear owner, deadline, or priority assigned"),
    ("😩", "Teams lose 30%+ of meeting value to poor follow-up"),
]

for i, (icon, problem) in enumerate(problems):
    y = Inches(1.35 + i * 1.05)
    card = box(slide, Inches(0.5), y, Inches(12.3), Inches(0.85), CARD_BG)
    # icon circle
    box(slide, Inches(0.6), Inches(y/Emu(914400) + 0.08), Inches(0.7), Inches(0.68), ACCENT2)
    txt(slide, icon, Inches(0.6), Inches(y/Emu(914400) + 0.06), Inches(0.7), Inches(0.72),
        size=22, align=PP_ALIGN.CENTER)
    txt(slide, problem, Inches(1.45), Inches(y/Emu(914400) + 0.12), Inches(11), Inches(0.55),
        size=19, color=WHITE)

box(slide, Inches(0), Inches(6.8), W, Inches(0.07), ACCENT2)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Our Solution
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
bg(slide)
box(slide, Inches(0), Inches(0), W, Inches(0.07), ACCENT3)

txt(slide, "✨  Our Solution", Inches(0.5), Inches(0.2), Inches(10), Inches(0.8),
    size=38, bold=True, color=ACCENT3)

txt(slide, "Paste your meeting notes → AI extracts every task → One-click push to Monday.com",
    Inches(0.5), Inches(1.0), Inches(12.3), Inches(0.55),
    size=18, italic=True, color=LIGHT_GRAY)

# Flow diagram boxes
steps = [
    ("📋", "Paste\nMeeting Notes", ACCENT2),
    ("🤖", "Local AI\nExtracts Tasks", ACCENT),
    ("✏️", "Review &\nEdit Tasks", ACCENT3),
    ("🚀", "Push to\nMonday.com", YELLOW),
]
box_w, box_h = Inches(2.6), Inches(2.8)
start_x = Inches(0.5)
for i, (icon, label, color) in enumerate(steps):
    x = start_x + i * Inches(3.2)
    y = Inches(1.9)
    # shadow
    box(slide, x + Inches(0.07), y + Inches(0.07), box_w, box_h, DARK_BG)
    card = box(slide, x, y, box_w, box_h, CARD_BG)
    # top color strip
    box(slide, x, y, box_w, Inches(0.12), color)
    txt(slide, icon, x, y + Inches(0.3), box_w, Inches(0.9),
        size=38, align=PP_ALIGN.CENTER)
    txt(slide, label, x, y + Inches(1.2), box_w, Inches(1.4),
        size=17, bold=True, color=color, align=PP_ALIGN.CENTER)
    # arrow between boxes
    if i < 3:
        txt(slide, "➜", x + box_w, y + Inches(1.0), Inches(0.7), Inches(0.7),
            size=28, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

box(slide, Inches(0), Inches(6.8), W, Inches(0.07), ACCENT3)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Key Features
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
bg(slide)
box(slide, Inches(0), Inches(0), W, Inches(0.07), ACCENT)

txt(slide, "⚡  Key Features", Inches(0.5), Inches(0.2), Inches(10), Inches(0.8),
    size=38, bold=True, color=ACCENT)

features = [
    ("🔒", "100% Private",          "Runs on your machine with Ollama — data never leaves.", ACCENT2),
    ("⚡", "3–7s Extraction",       "Optimised pipeline with smart caching for instant repeats.", ACCENT),
    ("📄", "Any Format",            "Markdown, bullets, plain text, chat transcripts & more.", ACCENT3),
    ("✏️", "Editable Preview",      "Review, modify, or delete tasks before pushing.", YELLOW),
    ("🔗", "Monday.com Native",     "One click sends tasks with assignees, dates & priorities.", ACCENT2),
    ("🎨", "Beautiful UI",          "Material-UI with dark theme and real-time updates.", ACCENT),
]

cols = 2
col_w = Inches(6.3)
for i, (icon, title, desc, color) in enumerate(features):
    col = i % cols
    row = i // cols
    x = Inches(0.4) + col * Inches(6.7)
    y = Inches(1.3) + row * Inches(1.75)
    card = box(slide, x, y, col_w, Inches(1.55), CARD_BG)
    box(slide, x, y, Inches(0.1), Inches(1.55), color)
    txt(slide, icon + "  " + title, x + Inches(0.25), y + Inches(0.12), col_w - Inches(0.3), Inches(0.55),
        size=18, bold=True, color=color)
    txt(slide, desc, x + Inches(0.25), y + Inches(0.65), col_w - Inches(0.35), Inches(0.75),
        size=13, color=LIGHT_GRAY)

box(slide, Inches(0), Inches(6.8), W, Inches(0.07), ACCENT)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Architecture
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
bg(slide)
box(slide, Inches(0), Inches(0), W, Inches(0.07), ACCENT2)

txt(slide, "🏗️  Architecture", Inches(0.5), Inches(0.2), Inches(10), Inches(0.8),
    size=38, bold=True, color=ACCENT2)

# arch boxes
arch = [
    (Inches(0.4),  Inches(1.5),  "🖥️ React Frontend\nTypeScript + MUI",         ACCENT2),
    (Inches(4.5),  Inches(1.5),  "⚙️ FastAPI Backend\nPython + Pydantic",         ACCENT),
    (Inches(8.6),  Inches(1.0),  "🤖 Ollama\nLocal LLM",                          ACCENT3),
    (Inches(8.6),  Inches(3.5),  "📌 Monday.com\nGraphQL API",                     YELLOW),
]

for (x, y, label, color) in arch:
    box(slide, x, y, Inches(3.5), Inches(1.6), CARD_BG)
    box(slide, x, y, Inches(3.5), Inches(0.12), color)
    txt(slide, label, x, y + Inches(0.2), Inches(3.5), Inches(1.3),
        size=17, bold=True, color=color, align=PP_ALIGN.CENTER)

# connector arrows text
txt(slide, "HTTP/REST ➜", Inches(3.95), Inches(2.1), Inches(1.5), Inches(0.4),
    size=12, italic=True, color=LIGHT_GRAY)
txt(slide, "➜ REST API\n(localhost:11434)", Inches(8.15), Inches(2.3), Inches(1.7), Inches(0.6),
    size=11, italic=True, color=LIGHT_GRAY)
txt(slide, "➜ GraphQL\n(Cloud)", Inches(8.15), Inches(4.3), Inches(1.5), Inches(0.5),
    size=11, italic=True, color=LIGHT_GRAY)

# tech stack list
txt(slide, "Tech Stack", Inches(0.4), Inches(3.4), Inches(7), Inches(0.55),
    size=20, bold=True, color=WHITE)

stack_items = [
    ("Frontend",     "React 18 · TypeScript · Material-UI · Vite",            ACCENT2),
    ("Backend",      "FastAPI · Pydantic · httpx · uvicorn · Python 3.11+",   ACCENT),
    ("AI Layer",     "Ollama · llama3.2:1b · Local inference (no cloud)",      ACCENT3),
    ("Integration",  "Monday.com GraphQL API · Token auth · Real-time push",   YELLOW),
]

for i, (cat, detail, color) in enumerate(stack_items):
    y = Inches(4.1) + i * Inches(0.65)
    box(slide, Inches(0.4), y, Inches(0.18), Inches(0.42), color)
    txt(slide, f"{cat}:  {detail}", Inches(0.7), y, Inches(7.5), Inches(0.45),
        size=13, color=LIGHT_GRAY)

box(slide, Inches(0), Inches(6.8), W, Inches(0.07), ACCENT2)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Demo / Screenshot
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
bg(slide)
box(slide, Inches(0), Inches(0), W, Inches(0.07), ACCENT3)

txt(slide, "🎬  How It Works", Inches(0.5), Inches(0.2), Inches(10), Inches(0.75),
    size=38, bold=True, color=ACCENT3)

# Mock UI screenshot using shapes
# Browser chrome
browser = box(slide, Inches(0.4), Inches(1.1), Inches(8.5), Inches(5.8), CARD_BG)
box(slide, Inches(0.4), Inches(1.1), Inches(8.5), Inches(0.45), RGBColor(0x1E, 0x3A, 0x50))
# browser dots
for di, dcol in enumerate([RGBColor(0xFF,0x5F,0x57), RGBColor(0xFF,0xBD,0x2E), RGBColor(0x28,0xCA,0x41)]):
    cx = Inches(0.65 + di*0.25)
    box(slide, cx, Inches(1.22), Inches(0.12), Inches(0.18), dcol)

# URL bar
box(slide, Inches(1.2), Inches(1.17), Inches(5.5), Inches(0.28), RGBColor(0x0D,0x1B,0x2A))
txt(slide, "localhost:5173  — MeetingToTask", Inches(1.22), Inches(1.14), Inches(5.5), Inches(0.3),
    size=9, color=LIGHT_GRAY)

# Mock app content
txt(slide, "MeetingToTask", Inches(0.6), Inches(1.65), Inches(4), Inches(0.5),
    size=20, bold=True, color=ACCENT)

# Input box mock
box(slide, Inches(0.6), Inches(2.2), Inches(3.9), Inches(2.8), RGBColor(0x0D,0x1B,0x2A))
txt(slide, "Meeting Notes", Inches(0.65), Inches(2.22), Inches(3), Inches(0.35),
    size=11, color=ACCENT3)
sample_notes = ("Q4 kickoff — April 22\n"
    "• Alice to finalize DB schema by Friday\n"
    "• Bob: set up CI/CD pipeline ASAP (high)\n"
    "• Carol review design mockups — due Mon\n"
    "• Everyone: update Jira tickets today")
txt(slide, sample_notes, Inches(0.65), Inches(2.6), Inches(3.8), Inches(2.2),
    size=10, color=LIGHT_GRAY)

# Extract button
btn = box(slide, Inches(0.6), Inches(5.15), Inches(3.9), Inches(0.42), ACCENT)
txt(slide, "⚡  Extract Tasks", Inches(0.6), Inches(5.14), Inches(3.9), Inches(0.44),
    size=14, bold=True, color=DARK_BG, align=PP_ALIGN.CENTER)

# Task table mock
box(slide, Inches(4.6), Inches(2.2), Inches(4.1), Inches(3.37), RGBColor(0x0D,0x1B,0x2A))
txt(slide, "Extracted Tasks (4)", Inches(4.65), Inches(2.22), Inches(4), Inches(0.35),
    size=11, color=ACCENT3)

tasks_mock = [
    ("Finalize DB schema",      "Alice",  "Fri",  "Medium", ACCENT3),
    ("Set up CI/CD pipeline",   "Bob",    "ASAP", "High",   RGBColor(0xFF,0x5F,0x57)),
    ("Review design mockups",   "Carol",  "Mon",  "Medium", ACCENT),
    ("Update Jira tickets",     "All",    "Today","Low",    ACCENT2),
]

for i, (task, who, due, pri, col) in enumerate(tasks_mock):
    y = Inches(2.65 + i * 0.72)
    box(slide, Inches(4.6), Inches(y/Emu(914400)), Inches(4.1), Inches(0.6), CARD_BG)
    box(slide, Inches(4.6), Inches(y/Emu(914400)), Inches(0.07), Inches(0.6), col)
    txt(slide, task, Inches(4.75), Inches(y/Emu(914400) + 0.04), Inches(2.2), Inches(0.3),
        size=10, bold=True, color=WHITE)
    txt(slide, f"👤 {who}  📅 {due}", Inches(4.75), Inches(y/Emu(914400) + 0.32), Inches(2.2), Inches(0.25),
        size=9, color=LIGHT_GRAY)
    box(slide, Inches(6.85), Inches(y/Emu(914400) + 0.08), Inches(0.9), Inches(0.28), col)
    txt(slide, pri, Inches(6.85), Inches(y/Emu(914400) + 0.06), Inches(0.9), Inches(0.32),
        size=9, bold=True, color=DARK_BG, align=PP_ALIGN.CENTER)

# Push button
box(slide, Inches(4.6), Inches(5.62), Inches(4.1), Inches(0.38), YELLOW)
txt(slide, "🚀  Push to Monday.com", Inches(4.6), Inches(5.61), Inches(4.1), Inches(0.4),
    size=13, bold=True, color=DARK_BG, align=PP_ALIGN.CENTER)

# Right side callouts
callouts = [
    ("🤖", "Local AI\n(Ollama)", Inches(9.2), Inches(1.5)),
    ("⚡", "3–7s\nExtraction", Inches(10.8), Inches(2.5)),
    ("🔒", "100%\nPrivate", Inches(9.2), Inches(3.7)),
    ("📌", "Monday.com\nIntegration", Inches(10.8), Inches(4.7)),
]
for icon, label, cx, cy in callouts:
    box(slide, cx, cy, Inches(1.5), Inches(1.2), CARD_BG)
    box(slide, cx, cy, Inches(1.5), Inches(0.08), ACCENT)
    txt(slide, icon, cx, cy + Inches(0.1), Inches(1.5), Inches(0.55),
        size=24, align=PP_ALIGN.CENTER)
    txt(slide, label, cx, cy + Inches(0.6), Inches(1.5), Inches(0.55),
        size=10, bold=True, color=ACCENT3, align=PP_ALIGN.CENTER)

box(slide, Inches(0), Inches(6.8), W, Inches(0.07), ACCENT3)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Performance & Stats
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
bg(slide)
box(slide, Inches(0), Inches(0), W, Inches(0.07), YELLOW)

txt(slide, "📊  Performance & Impact", Inches(0.5), Inches(0.2), Inches(10), Inches(0.75),
    size=38, bold=True, color=YELLOW)

# Big stat cards
stats = [
    ("3–7s",  "Average Task\nExtraction Time",   ACCENT),
    ("⚡ Instant", "Cached\nResults",             ACCENT3),
    ("🔒 0",  "Data Sent\nto Cloud",              ACCENT2),
    ("8+",    "Meeting Formats\nSupported",        YELLOW),
]

for i, (val, label, color) in enumerate(stats):
    x = Inches(0.4 + i * 3.2)
    y = Inches(1.4)
    box(slide, x, y, Inches(2.9), Inches(2.5), CARD_BG)
    box(slide, x, y, Inches(2.9), Inches(0.14), color)
    txt(slide, val, x, y + Inches(0.3), Inches(2.9), Inches(1.1),
        size=40, bold=True, color=color, align=PP_ALIGN.CENTER)
    txt(slide, label, x, y + Inches(1.35), Inches(2.9), Inches(1.0),
        size=14, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# Performance breakdown bar chart (visual)
txt(slide, "Speed improvements vs baseline", Inches(0.4), Inches(4.1), Inches(8), Inches(0.45),
    size=16, bold=True, color=WHITE)

bars = [
    ("Before (v1)", 25, RGBColor(0x88,0x44,0x44)),
    ("After (v2)",  6,  ACCENT3),
    ("Cached",      0.2, ACCENT),
]

for i, (label, sec, color) in enumerate(bars):
    y = Inches(4.7 + i * 0.9)
    txt(slide, label, Inches(0.4), y, Inches(1.8), Inches(0.55), size=13, color=LIGHT_GRAY)
    bar_w = Inches(min(sec / 30 * 9.5, 9.5))
    box(slide, Inches(2.3), y + Inches(0.08), bar_w, Inches(0.38), color)
    txt(slide, f"{sec}s", Inches(2.35 + min(sec/30*9.5, 9.5)), y, Inches(1), Inches(0.5),
        size=13, bold=True, color=WHITE)

box(slide, Inches(0), Inches(6.8), W, Inches(0.07), YELLOW)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Why We Win (USP)
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
bg(slide)
box(slide, Inches(0), Inches(0), W, Inches(0.07), ACCENT2)

txt(slide, "🎯  Why MeetingToTask Wins", Inches(0.5), Inches(0.2), Inches(10), Inches(0.75),
    size=38, bold=True, color=ACCENT2)

# Comparison table header
headers = ["Feature", "MeetingToTask", "Manual Entry", "Generic AI Tools"]
col_widths = [Inches(3.0), Inches(2.6), Inches(2.6), Inches(2.6)]
col_offsets = [Inches(0.4), Inches(3.5), Inches(6.2), Inches(8.85)]

for j, (h, w, x) in enumerate(zip(headers, col_widths, col_offsets)):
    hdr_color = ACCENT if j == 1 else CARD_BG
    box(slide, x, Inches(1.2), w, Inches(0.55), hdr_color)
    txt(slide, h, x, Inches(1.21), w, Inches(0.53),
        size=14, bold=True, color=WHITE if j != 1 else DARK_BG, align=PP_ALIGN.CENTER)

rows = [
    ("🔒 Data Privacy",     "✅ 100% Local",  "✅ Local",      "❌ Cloud"),
    ("⚡ Speed",             "✅ 3–7 seconds", "❌ 5–30 min",   "⚠️ 10–20s"),
    ("🔗 Monday.com Push",  "✅ One-click",   "✅ Manual",     "❌ None"),
    ("📄 Multi-format",     "✅ 8+ formats",  "⚠️ Tedious",    "⚠️ Limited"),
    ("💰 Cost",             "✅ Free / OSS",  "✅ Free",       "❌ Paid API"),
    ("🎯 Task Intelligence","✅ Assignee+Date","❌ None",        "⚠️ Basic"),
]

for i, row in enumerate(rows):
    for j, (cell, w, x) in enumerate(zip(row, col_widths, col_offsets)):
        row_bg = CARD_BG if i % 2 == 0 else RGBColor(0x18, 0x30, 0x46)
        if j == 1: row_bg = RGBColor(0x00, 0x3A, 0x4F)
        box(slide, x, Inches(1.8 + i*0.78), w, Inches(0.72), row_bg)
        c = ACCENT3 if j == 1 else (WHITE if cell.startswith("✅") else (YELLOW if cell.startswith("⚠️") else RGBColor(0xFF,0x5F,0x57)))
        txt(slide, cell, x, Inches(1.81 + i*0.78), w, Inches(0.68),
            size=12, color=c, align=PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT)

box(slide, Inches(0), Inches(6.8), W, Inches(0.07), ACCENT2)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Closing / CTA
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
bg(slide)

# Full gradient background
grad2 = gradient_image(960, 540, (0,40,80), (13,27,42))
add_pil_image(slide, grad2, Inches(0), Inches(0), W, H)

box(slide, Inches(0), Inches(0), W, Inches(0.08), ACCENT)
box(slide, Inches(0), Inches(7.42), W, Inches(0.08), ACCENT)

# big center text
txt(slide, "🚀", Inches(0), Inches(0.8), W, Inches(1.2),
    size=60, align=PP_ALIGN.CENTER)

txt(slide, "MeetingToTask", Inches(0), Inches(1.9), W, Inches(1.2),
    size=60, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

txt(slide, "From Meeting to Monday — in Seconds.",
    Inches(0), Inches(3.0), W, Inches(0.75),
    size=26, color=WHITE, align=PP_ALIGN.CENTER)

txt(slide, "Open Source  ·  Local AI  ·  Zero Friction",
    Inches(0), Inches(3.7), W, Inches(0.6),
    size=18, color=ACCENT3, align=PP_ALIGN.CENTER)

# CTA button
box(slide, Inches(4.9), Inches(4.6), Inches(3.5), Inches(0.7), ACCENT)
txt(slide, "Try It Now →", Inches(4.9), Inches(4.59), Inches(3.5), Inches(0.72),
    size=22, bold=True, color=DARK_BG, align=PP_ALIGN.CENTER)

txt(slide, "github.com/your-team/meetingtotask",
    Inches(0), Inches(5.6), W, Inches(0.5),
    size=14, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

txt(slide, "Built at Hackathon 2026  ·  April 22, 2026",
    Inches(0), Inches(6.5), W, Inches(0.5),
    size=12, italic=True, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# ── Save ───────────────────────────────────────────────────────────────────────
out = "/tmp/meetingtotask/MeetingToTask_Hackathon.pptx"
prs.save(out)
print(f"✅  Saved: {out}")
print(f"   Slides: {len(prs.slides)}")
